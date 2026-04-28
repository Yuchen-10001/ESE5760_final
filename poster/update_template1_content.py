from __future__ import annotations

from pathlib import Path
from shutil import copyfile

import matplotlib.pyplot as plt
import matplotlib.image as mpimg
from matplotlib.patches import FancyBboxPatch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt


ROOT = Path(__file__).resolve().parent.parent
POSTER = ROOT / "poster"
TEMPLATE = POSTER / "Template1A_36x48.pptx"
BACKUP = POSTER / "Template1A_36x48.backup.pptx"
ASSETS = POSTER / "template1_assets"
PLOTS = ROOT / "experiments" / "plots"

BLUE = RGBColor(0x10, 0x2E, 0x66)
BLACK = RGBColor(0x11, 0x11, 0x11)
ORANGE = RGBColor(0xD8, 0x6A, 0x00)


def build_assets() -> dict[str, Path]:
    ASSETS.mkdir(parents=True, exist_ok=True)

    experiment = ASSETS / "experiment_matrix.png"
    run_summary = ASSETS / "run_summary.png"
    results_panel = ASSETS / "main_results_panel.png"
    tsv_panel = ASSETS / "tsv_results_panel.png"

    build_experiment_matrix(experiment)
    build_run_summary(run_summary)
    build_main_results_panel(results_panel)
    build_tsv_panel(tsv_panel)

    return {
        "experiment": experiment,
        "run_summary": run_summary,
        "results_panel": results_panel,
        "tsv_panel": tsv_panel,
    }


def build_experiment_matrix(out_path: Path) -> None:
    fig = plt.figure(figsize=(7.4, 5.2), dpi=180, facecolor="white")
    ax = fig.add_axes([0.03, 0.04, 0.94, 0.92])
    ax.axis("off")

    def box(x, y, w, h, title, lines, fc, ec="#5E7FB9"):
        patch = FancyBboxPatch(
            (x, y),
            w,
            h,
            boxstyle="round,pad=0.02,rounding_size=0.03",
            facecolor=fc,
            edgecolor=ec,
            linewidth=1.8,
        )
        ax.add_patch(patch)
        ax.text(x + w / 2, y + h - 0.08, title, ha="center", va="top", fontsize=13, fontweight="bold", color="#102E66")
        ypos = y + h - 0.16
        for line in lines:
            ax.text(x + 0.03, ypos, f"- {line}", ha="left", va="top", fontsize=10.5, color="#222222")
            ypos -= 0.11

    ax.text(0.5, 0.96, "Proposal-Aligned Experiment Structure", ha="center", va="top", fontsize=16, fontweight="bold", color="#102E66")
    ax.text(0.5, 0.90, "Poster scope keeps the fixed-capacity comparison, node scaling, and TSV robustness checks.", ha="center", va="top", fontsize=10.5, color="#39599A")

    box(
        0.03,
        0.48,
        0.28,
        0.32,
        "Normalized Baseline",
        [
            "2 MB cache, 256-bit word width, 1-way associativity",
            "LOP roadmap, 350 K operating point",
            "DESTINY optimization target = WriteEDP",
        ],
        "#EAF1FF",
    )
    box(
        0.36,
        0.48,
        0.28,
        0.32,
        "Experiment 1",
        [
            "2D SRAM, 3D SRAM, 2D eDRAM, 3D eDRAM",
            "65 nm, 45 nm, 32 nm, 22 nm",
            "16 runs for node-scaling trends",
        ],
        "#F4F8E8",
    )
    box(
        0.69,
        0.48,
        0.28,
        0.32,
        "Experiment 2",
        [
            "3D SRAM fixed at 32 nm",
            "Sweep LocalTSV, GlobalTSV, TSVRedundancy",
            "6 TSV runs + reused baseline",
        ],
        "#FFF1E5",
    )

    ax.annotate("", xy=(0.36, 0.64), xytext=(0.31, 0.64), arrowprops=dict(arrowstyle="->", lw=2, color="#5E7FB9"))
    ax.annotate("", xy=(0.69, 0.64), xytext=(0.64, 0.64), arrowprops=dict(arrowstyle="->", lw=2, color="#5E7FB9"))

    ax.text(0.5, 0.29, "Extracted metrics for every run", ha="center", va="center", fontsize=13, fontweight="bold", color="#102E66")
    metrics = [
        "Total area (mm^2)",
        "Read latency (ns)",
        "Write latency (ns)",
        "Write dynamic energy (nJ/access)",
        "Leakage power (mW)",
    ]
    xs = [0.10, 0.30, 0.50, 0.70, 0.90]
    for x, metric in zip(xs, metrics):
        circ = plt.Circle((x, 0.13), 0.07, facecolor="#EAF1FF", edgecolor="#5E7FB9", linewidth=1.5)
        ax.add_patch(circ)
        ax.text(x, 0.13, metric, ha="center", va="center", fontsize=8.5, color="#102E66", wrap=True)

    fig.savefig(out_path, bbox_inches="tight")
    plt.close(fig)


def build_run_summary(out_path: Path) -> None:
    fig = plt.figure(figsize=(4.4, 4.2), dpi=180, facecolor="white")
    ax = fig.add_axes([0, 0, 1, 1])
    ax.axis("off")

    circ1 = plt.Circle((0.33, 0.56), 0.27, facecolor="#EAF1FF", edgecolor="#5E7FB9", linewidth=2)
    circ2 = plt.Circle((0.67, 0.56), 0.23, facecolor="#FFF1E5", edgecolor="#D86A00", linewidth=2)
    ax.add_patch(circ1)
    ax.add_patch(circ2)

    ax.text(0.33, 0.63, "16", ha="center", va="center", fontsize=28, fontweight="bold", color="#102E66")
    ax.text(0.33, 0.49, "node-scaling runs", ha="center", va="center", fontsize=11, color="#102E66")
    ax.text(0.67, 0.63, "6", ha="center", va="center", fontsize=24, fontweight="bold", color="#B35600")
    ax.text(0.67, 0.49, "TSV sweeps", ha="center", va="center", fontsize=11, color="#B35600")
    ax.text(0.50, 0.18, "22 total simulation points drive the poster conclusions.", ha="center", va="center", fontsize=11, color="#222222")

    fig.savefig(out_path, bbox_inches="tight")
    plt.close(fig)


def build_main_results_panel(out_path: Path) -> None:
    fig = plt.figure(figsize=(10.5, 8.2), dpi=180, facecolor="white")
    gs = fig.add_gridspec(2, 2, left=0.04, right=0.98, top=0.92, bottom=0.08, hspace=0.28, wspace=0.18)

    panels = [
        (PLOTS / "fig3_total_area_vs_node.png", "Area scaling"),
        (PLOTS / "fig1_read_latency_vs_node.png", "Read latency"),
        (PLOTS / "fig4_leakage_power_vs_node.png", "Leakage behavior"),
    ]
    slots = [(0, 0), (0, 1), (1, 0)]
    for (path, title), (r, c) in zip(panels, slots):
        ax = fig.add_subplot(gs[r, c])
        ax.imshow(mpimg.imread(path))
        ax.set_axis_off()
        ax.set_title(title, fontsize=12, fontweight="bold", color="#102E66", pad=6)

    ax = fig.add_subplot(gs[1, 1])
    ax.axis("off")
    ax.text(0.5, 0.95, "Representative 32 nm comparison", ha="center", va="top", fontsize=12, fontweight="bold", color="#102E66")
    table_rows = [
        ["2D SRAM", "4.262", "1.061", "0.123", "234.6"],
        ["3D SRAM", "2.132", "0.633", "0.089", "234.6"],
        ["2D eDRAM", "1.440", "0.757", "0.066", "30.0"],
        ["3D eDRAM", "1.068", "0.472", "0.061", "50.1"],
    ]
    table = ax.table(
        cellText=table_rows,
        colLabels=["Tech", "Area", "Read", "Write E.", "Leakage"],
        cellLoc="center",
        colLoc="center",
        bbox=[0.02, 0.28, 0.96, 0.55],
    )
    table.auto_set_font_size(False)
    table.set_fontsize(9.5)
    for (r, c), cell in table.get_celld().items():
        cell.set_edgecolor("#A6BEE8")
        if r == 0:
            cell.set_facecolor("#DDE8FF")
            cell.set_text_props(weight="bold", color="#102E66")
        else:
            cell.set_facecolor("#F8FAFE" if r % 2 else "#EEF4FF")
    for key in [(4, 1), (4, 2), (4, 3), (3, 4)]:
        table[key].set_facecolor("#DDF3E4")
        table[key].set_text_props(weight="bold", color="#0B6E3E")
    bullets = [
        "3D eDRAM leads active metrics at 32 nm.",
        "2D eDRAM is the lowest-leakage standby option.",
        "3D SRAM is the balanced middle ground when refresh is undesirable.",
    ]
    y = 0.20
    for bullet in bullets:
        ax.text(0.04, y, f"- {bullet}", ha="left", va="top", fontsize=9.5, color="#222222")
        y -= 0.08

    fig.suptitle("Main Analytical Results", fontsize=16, fontweight="bold", color="#102E66")
    fig.savefig(out_path, bbox_inches="tight")
    plt.close(fig)


def build_tsv_panel(out_path: Path) -> None:
    fig = plt.figure(figsize=(9.2, 6.3), dpi=180, facecolor="white")
    gs = fig.add_gridspec(2, 1, left=0.05, right=0.98, top=0.90, bottom=0.10, hspace=0.18)
    for idx, (path, title) in enumerate(
        [
            (PLOTS / "fig5_tsv_sensitivity_area.png", "Area sensitivity"),
            (PLOTS / "fig6_tsv_sensitivity_latency.png", "Latency sensitivity"),
        ]
    ):
        ax = fig.add_subplot(gs[idx, 0])
        ax.imshow(mpimg.imread(path))
        ax.set_axis_off()
        ax.set_title(title, fontsize=11.5, fontweight="bold", color="#102E66", pad=6)
    fig.suptitle("TSV Sensitivity at 3D SRAM, 32 nm", fontsize=15, fontweight="bold", color="#102E66")
    fig.savefig(out_path, bbox_inches="tight")
    plt.close(fig)


def walk_shapes(shapes):
    for shape in shapes:
        yield shape
        if hasattr(shape, "shapes"):
            yield from walk_shapes(shape.shapes)


def find_shape(slide, name: str):
    for shape in walk_shapes(slide.shapes):
        if shape.name == name:
            return shape
    raise KeyError(name)


def set_text(shape, lines, *, size, bold=False, color=BLUE, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_picture(slide, path: Path, left: int, top: int, width: int, height: int):
    slide.shapes.add_picture(str(path), Emu(left), Emu(top), width=Emu(width), height=Emu(height))


def update_template():
    assets = build_assets()
    if not BACKUP.exists():
        copyfile(TEMPLATE, BACKUP)

    prs = Presentation(str(TEMPLATE))
    slide = prs.slides[0]

    set_text(
        find_shape(slide, "TextBox 56"),
        ["Comparative Analysis of 2D/3D SRAM and eDRAM Across Process Nodes Using DESTINY"],
        size=26,
        bold=True,
        color=BLUE,
        align=PP_ALIGN.CENTER,
    )
    set_text(
        find_shape(slide, "TextBox 57"),
        ["Tiebing Tang  •  Yishu Wang  •  Zhengyu Liu"],
        size=18,
        color=BLUE,
        align=PP_ALIGN.CENTER,
    )
    set_text(
        find_shape(slide, "TextBox 58"),
        ["ESE5760 Final Project  •  Proposal-guided poster scope  •  Conclusions derived from Doc/Analysis_Full.md"],
        size=14,
        color=BLUE,
        align=PP_ALIGN.CENTER,
    )

    set_text(find_shape(slide, "TextBox 34"), ["Motivation & Scope"], size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    set_text(
        find_shape(slide, "TextBox 60"),
        [
            "The proposal framed this project around one question: which memory trends are robust, and which are highly assumption-dependent in DESTINY?",
            "",
            "We therefore narrowed the poster to the most stable and interpretable comparisons:",
            "• fixed-capacity cross-technology comparison,",
            "• node-scaling trends across 65/45/32/22 nm,",
            "• 2D/3D benefit and TSV-sensitive robustness.",
            "",
            "The poster intentionally favors a few strong analytical results over exhaustive coverage of every simulation output.",
        ],
        size=15,
        color=BLACK,
    )

    set_text(find_shape(slide, "TextBox 35"), ["Methodology & Experiment Matrix"], size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    set_text(
        find_shape(slide, "TextBox 61"),
        [
            "All runs share a normalized baseline: 2 MB capacity, 256-bit word width, 1-way associativity, LOP roadmap, 350 K, and WriteEDP optimization.",
            "",
            "Experiment 1 compares 2D SRAM, 3D SRAM, 2D eDRAM, and 3D eDRAM across four nodes (16 runs).",
            "",
            "Experiment 2 fixes 3D SRAM at 32 nm and sweeps LocalTSVProjection, GlobalTSVProjection, and TSVRedundancy (6 runs).",
            "",
            "Because DESTINY optimizes for WriteEDP, some anomalies are organizational effects rather than pure device limits.",
        ],
        size=14,
        color=BLACK,
    )
    set_text(
        find_shape(slide, "TextBox 62"),
        ["Proposal filter: we keep the stable poster-ready experiments and defer broader sensitivity work to the full writeup."],
        size=11,
        color=BLACK,
        align=PP_ALIGN.CENTER,
    )
    set_text(
        find_shape(slide, "TextBox 59"),
        ["22 runs total: enough for comparative trends, small enough to explain clearly."],
        size=11,
        color=BLACK,
        align=PP_ALIGN.CENTER,
    )

    set_text(find_shape(slide, "TextBox 63"), ["Cross-Technology Results"], size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    set_text(
        find_shape(slide, "TextBox 65"),
        [
            "At the poster level, three trends are the most defensible:",
            "",
            "1. eDRAM scales more aggressively than SRAM, and 3D eDRAM becomes the smallest design by 22 nm (0.550 mm^2).",
            "2. 3D SRAM consistently saves roughly 42-50% area over 2D SRAM, so the stacking benefit is robust for SRAM.",
            "3. At 32 nm, 3D eDRAM wins the active metrics (area, read latency, write energy), while 2D eDRAM remains the best leakage option.",
        ],
        size=14,
        color=BLACK,
    )

    set_text(find_shape(slide, "TextBox 64"), ["Why 32 nm Is the Best Comparison Point"], size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    set_text(
        find_shape(slide, "TextBox 66"),
        [
            "32 nm is the cleanest fixed-capacity comparison node: it avoids the extreme 65 nm H-tree penalty while still preserving strong technology separation.",
            "",
            "The 32 nm scorecard shows a useful split:",
            "• 3D eDRAM is best when active performance matters.",
            "• 2D eDRAM is best when standby leakage dominates.",
            "• 3D SRAM is a balanced compromise if refresh complexity is undesirable.",
        ],
        size=13,
        color=BLACK,
    )

    set_text(find_shape(slide, "TextBox 68"), ["Interpretation of the Anomalies"], size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    set_text(
        find_shape(slide, "TextBox 69"),
        [
            "Two non-monotonic SRAM behaviors need interpretation rather than blind reporting.",
            "",
            "The 65 nm SRAM latency spike is caused by large-array H-tree delay under a WriteEDP-selected organization; it is not a fundamental SRAM speed limit.",
            "",
            "The 32 nm SRAM leakage spike (234.6 mW for both 2D and 3D SRAM) indicates that DESTINY converged to a high-parallelism organization that favors WriteEDP at the cost of standby power.",
            "",
            "This distinction matters because the proposal explicitly asked for robust trends versus assumption-sensitive conclusions. The area ordering and eDRAM leakage advantage are robust. The exact SRAM anomaly magnitudes depend on the optimizer target.",
        ],
        size=13,
        color=BLACK,
    )

    set_text(find_shape(slide, "TextBox 70"), ["TSV Sensitivity & 3D Robustness"], size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    set_text(
        find_shape(slide, "TextBox 74"),
        [
            "The TSV experiment answers the proposal's third research question directly.",
            "",
            "LocalTSVProjection is the only first-order area knob, and even there the maximum penalty is only +5.5%. GlobalTSVProjection and TSVRedundancy remain below 1% area impact.",
            "",
            "Read latency is unchanged across all TSV sweeps, so the critical path at this 2 MB scale is not TSV-limited.",
            "",
            "Practical implication: conservative TSV assumptions are acceptable for this class of cache without materially changing the system-level conclusion.",
        ],
        size=13,
        color=BLACK,
    )

    set_text(find_shape(slide, "TextBox 71"), ["Key Conclusions"], size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    set_text(
        find_shape(slide, "TextBox 75"),
        [
            "1. Node scaling helps every technology, but eDRAM gains more from scaling than SRAM.",
            "",
            "2. 3D stacking gives a clear and repeatable area benefit for SRAM; for eDRAM it becomes beneficial only after TSV overhead stops dominating.",
            "",
            "3. At 32 nm, 3D eDRAM is the strongest active design point, while 2D eDRAM is the most attractive low-leakage option.",
            "",
            "4. The main 3D conclusion is robust even when TSV assumptions are made more conservative.",
        ],
        size=13,
        color=BLACK,
    )
    set_text(
        find_shape(slide, "TextBox 76"),
        ["Full writeup: Doc/Analysis_Full.md"],
        size=12,
        color=ORANGE,
        align=PP_ALIGN.CENTER,
    )

    set_text(find_shape(slide, "TextBox 72"), ["Limits & Design Implications"], size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    set_text(
        find_shape(slide, "TextBox 77"),
        [
            "DESTINY does not include eDRAM refresh energy, so the eDRAM advantage here is strongest on active metrics rather than full-system power.",
            "",
            "22 nm should be interpreted as a model projection, not silicon validation.",
            "",
            "All reported trends are under a single optimization target (WriteEDP). A ReadLatency objective would likely weaken the SRAM anomalies and strengthen some 3D latency gains.",
        ],
        size=12.5,
        color=BLACK,
    )
    set_text(find_shape(slide, "TextBox 78"), ["ROBUST TRENDS"], size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    set_text(find_shape(slide, "TextBox 79"), ["ASSUMPTION-SENSITIVE DETAILS"], size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    # Overlay analytical figures while keeping the template layout.
    add_picture(slide, assets["experiment"], 8585435, 14090168, 4737100, 4419600)
    add_picture(slide, assets["run_summary"], 3193735, 15539112, 2641600, 2425700)
    add_picture(slide, assets["results_panel"], 3969419, 25240911, 6080344, 5585771)
    add_picture(slide, assets["tsv_panel"], 33127419, 10831938, 7074056, 4707174)

    prs.save(str(TEMPLATE))


if __name__ == "__main__":
    update_template()
    print(TEMPLATE)

from __future__ import annotations

from pathlib import Path
from shutil import copyfile

import matplotlib.pyplot as plt
import qrcode
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt


ROOT = Path(__file__).resolve().parent.parent
POSTER_DIR = ROOT / "poster"
ASSET_DIR = POSTER_DIR / "generated_assets"
TEMPLATE = POSTER_DIR / "36x48_blue.pptx"
OUTPUT = POSTER_DIR / "36x48_blue_filled.pptx"

PLOTS = ROOT / "experiments" / "plots"
REPO_URL = "https://github.com/Yuchen-10001/ESE5760_final"
ANALYSIS_URL = f"{REPO_URL}/blob/main/Doc/Analysis_Full.md"

NAVY = RGBColor(0x01, 0x1F, 0x5B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x11, 0x11, 0x11)


def ensure_assets() -> dict[str, Path]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)

    repo_qr = ASSET_DIR / "repo_qr.png"
    analysis_qr = ASSET_DIR / "analysis_qr.png"
    scorecard = ASSET_DIR / "scorecard_32nm.png"

    build_qr(REPO_URL, repo_qr)
    build_qr(ANALYSIS_URL, analysis_qr)
    build_scorecard(scorecard)

    return {
        "repo_qr": repo_qr,
        "analysis_qr": analysis_qr,
        "scorecard": scorecard,
    }


def build_qr(url: str, out_path: Path) -> None:
    qr = qrcode.QRCode(border=1, box_size=12)
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    img.save(out_path)


def build_scorecard(out_path: Path) -> None:
    rows = [
        ["2D SRAM", "4.262", "1.061", "0.123", "234.6"],
        ["3D SRAM", "2.132", "0.633", "0.089", "234.6"],
        ["2D eDRAM", "1.440", "0.757", "0.066", "30.0"],
        ["3D eDRAM", "1.068", "0.472", "0.061", "50.1"],
    ]
    cols = ["Tech", "Area", "Read", "Write E.", "Leakage"]

    fig = plt.figure(figsize=(6.2, 6.2), dpi=200, facecolor="white")
    ax = fig.add_axes([0.05, 0.07, 0.90, 0.86])
    ax.axis("off")

    fig.text(
        0.5,
        0.95,
        "Representative Design Point: 32 nm",
        ha="center",
        va="top",
        fontsize=15,
        fontweight="bold",
        color="#011F5B",
    )
    fig.text(
        0.5,
        0.91,
        "Best active metrics: 3D eDRAM | Best leakage: 2D eDRAM",
        ha="center",
        va="top",
        fontsize=9.5,
        color="#274C9B",
    )

    table = ax.table(
        cellText=rows,
        colLabels=cols,
        cellLoc="center",
        colLoc="center",
        bbox=[0.0, 0.25, 1.0, 0.58],
    )
    table.auto_set_font_size(False)
    table.set_fontsize(10)

    for (r, c), cell in table.get_celld().items():
        if r == 0:
            cell.set_facecolor("#83ACFF")
            cell.set_text_props(color="#011F5B", weight="bold")
        else:
            cell.set_facecolor("#F6F8FC" if r % 2 else "#EAF0FF")
        cell.set_edgecolor("#9DBEFF")

    # Highlight best metrics in the comparison row set.
    highlight = {
        (4, 1),  # 3D eDRAM, area
        (4, 2),  # 3D eDRAM, read
        (4, 3),  # 3D eDRAM, write
        (3, 4),  # 2D eDRAM, leakage
    }
    for key in highlight:
        cell = table[key]
        cell.set_facecolor("#D9F2E3")
        cell.set_text_props(weight="bold", color="#0B6E3E")

    notes = [
        "3D SRAM halves area versus 2D SRAM at 32 nm (4.262 -> 2.132 mm^2).",
        "eDRAM keeps a large leakage advantage over SRAM (30.0 vs 234.6 mW).",
        "This 32 nm view is the proposal's fixed-capacity comparison panel.",
    ]
    y = 0.17
    for note in notes:
        fig.text(0.08, y, f"- {note}", ha="left", va="top", fontsize=9.5, color="#222222")
        y -= 0.05

    fig.savefig(out_path, bbox_inches="tight")
    plt.close(fig)


def walk_shapes(shapes):
    for shape in shapes:
        yield shape
        if hasattr(shape, "shapes"):
            yield from walk_shapes(shape.shapes)


def find_shape(slide, name: str):
    for shape in walk_shapes(slide.shapes):
        if shape.name == name:
            return shape
    raise KeyError(f"Shape {name!r} not found")


def clear_paragraphs(text_frame):
    while len(text_frame.paragraphs) > 1:
        p = text_frame.paragraphs[-1]
        p._element.getparent().remove(p._element)
    text_frame.paragraphs[0].clear()


def set_box_text(
    shape,
    lines: list[str],
    *,
    font_name: str,
    font_size: int,
    color: RGBColor,
    bold: bool = False,
    align=PP_ALIGN.CENTER,
    valign=MSO_ANCHOR.TOP,
) -> None:
    tf = shape.text_frame
    clear_paragraphs(tf)
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_body_textbox(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    lines: list[str],
    *,
    font_size: int = 20,
    color: RGBColor = NAVY,
    bullet: bool = False,
    bold_first: bool = False,
) -> None:
    shape = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if bullet:
            p.level = 0
        run = p.add_run()
        run.text = f"• {line}" if bullet else line
        run.font.name = "Arial MT Pro"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold_first and idx == 0


def add_picture(slide, path: Path, left: int, top: int, width: int, height: int) -> None:
    slide.shapes.add_picture(str(path), Emu(left), Emu(top), width=Emu(width), height=Emu(height))


def build_poster() -> Path:
    assets = ensure_assets()
    copyfile(TEMPLATE, OUTPUT)

    prs = Presentation(str(OUTPUT))
    slide = prs.slides[0]

    set_box_text(
        find_shape(slide, "TextBox 42"),
        ["Comparative Analysis of", "2D/3D SRAM and eDRAM", "Using DESTINY"],
        font_name="Arial MT Pro",
        font_size=34,
        color=WHITE,
        bold=False,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    set_box_text(
        find_shape(slide, "TextBox 43"),
        ["Tiebing Tang  •  Yishu Wang  •  Zhengyu Liu"],
        font_name="Arial MT Pro",
        font_size=22,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    set_box_text(
        find_shape(slide, "TextBox 32"),
        [
            "Proposal-guided takeaway:",
            "3D eDRAM is the strongest active design at 32 nm, 2D eDRAM keeps the best leakage,",
            "and TSV assumptions do not materially change the 2 MB cache conclusion.",
        ],
        font_name="Arial MT Pro",
        font_size=19,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    set_box_text(
        find_shape(slide, "TextBox 39"),
        [
            "Poster scope follows the proposal:",
            "fixed-capacity comparison, node scaling, and 2D/3D + TSV robustness.",
            "We emphasize stable trends and call out optimizer-driven anomalies separately.",
        ],
        font_name="Arial MT Pro",
        font_size=18,
        color=NAVY,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    set_box_text(
        find_shape(slide, "TextBox 35"),
        ["Project Links"],
        font_name="League Spartan",
        font_size=24,
        color=NAVY,
        align=PP_ALIGN.CENTER,
    )
    set_box_text(
        find_shape(slide, "TextBox 37"),
        ["GitHub repository"],
        font_name="Arial MT Pro",
        font_size=18,
        color=NAVY,
        align=PP_ALIGN.CENTER,
    )
    set_box_text(
        find_shape(slide, "TextBox 38"),
        ["Analysis + plots"],
        font_name="Arial MT Pro",
        font_size=18,
        color=NAVY,
        align=PP_ALIGN.CENTER,
    )
    set_box_text(
        find_shape(slide, "TextBox 13"),
        ["Full", "analysis"],
        font_name="Canva Sans",
        font_size=22,
        color=BLACK,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    set_box_text(
        find_shape(slide, "TextBox 25"),
        ["Code +", "data"],
        font_name="Canva Sans",
        font_size=18,
        color=BLACK,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )

    add_body_textbox(
        slide,
        left=950000,
        top=5900000,
        width=12800000,
        height=2300000,
        lines=[
            "3D eDRAM is best at 32 nm: 1.068 mm^2, 0.472 ns, 0.061 nJ/access.",
            "2D eDRAM gives the lowest leakage: 30.0 mW at 32 nm.",
            "3D SRAM cuts area by about 40-50% versus 2D SRAM across nodes.",
            "TSV sweeps change area by at most 5.5% and leave read latency unchanged.",
        ],
        font_size=17,
        color=BLACK,
        bullet=True,
    )
    add_body_textbox(
        slide,
        left=950000,
        top=9500000,
        width=12800000,
        height=1000000,
        lines=[
            "DESTINY simulator, ESE5760 course staff, and prior 3D-memory literature were used for",
            "baseline validation, sanity checks, and interpretation of anomaly cases.",
        ],
        font_size=15,
        color=BLACK,
    )
    add_body_textbox(
        slide,
        left=950000,
        top=12200000,
        width=13200000,
        height=8800000,
        lines=[
            "The proposal asked which technology trends remain stable when process node and 3D assumptions change in DESTINY.",
            "This poster narrows the final scope to the configurations that were both reliable and interpretable in simulation: 2D SRAM, 3D SRAM, 2D eDRAM, and 3D eDRAM.",
            "Our goal is not one perfect design point, but a robust comparison of what really changes with scaling, stacking, and TSV assumptions.",
        ],
        font_size=18,
        color=BLACK,
    )
    add_body_textbox(
        slide,
        left=950000,
        top=24150000,
        width=13200000,
        height=6500000,
        lines=[
            "Baseline: 2 MB cache, 256-bit word width, 1-way associativity, LOP roadmap, 350 K, and WriteEDP optimization.",
            "Experiment 1 sweeps four technologies across 65 nm, 45 nm, 32 nm, and 22 nm (16 runs total).",
            "Experiment 2 fixes 3D SRAM at 32 nm and sweeps LocalTSVProjection, GlobalTSVProjection, and TSVRedundancy.",
            "Metrics: total area, read latency, write latency, write dynamic energy, and leakage power.",
        ],
        font_size=18,
        color=BLACK,
    )

    # Results figures selected using the poster criteria in the proposal: 3-5 strong visuals.
    add_picture(slide, PLOTS / "fig3_total_area_vs_node.png", 13900000, 12150000, 7600000, 2350000)
    add_picture(slide, PLOTS / "fig1_read_latency_vs_node.png", 22300000, 12150000, 7600000, 2350000)
    add_picture(slide, PLOTS / "fig4_leakage_power_vs_node.png", 13900000, 21450000, 7900000, 4300000)
    add_picture(slide, PLOTS / "fig5_tsv_sensitivity_area.png", 22500000, 21450000, 7900000, 4300000)

    # Representative fixed-capacity comparison panel on the right.
    add_picture(slide, assets["scorecard"], 38095953, 718571, 5024943, 5043857)

    # QR overlays keep the template positions intact while replacing placeholder content.
    add_picture(slide, assets["repo_qr"], 40050384, 5380119, 2918266, 2910971)
    add_picture(slide, assets["analysis_qr"], 39187660, 27283970, 3987708, 3977738)

    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    out = build_poster()
    print(out)
